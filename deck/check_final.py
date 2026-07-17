"""Finalization gate for the deck that will actually be submitted.

Usage: python deck/check_final.py [path-to-pptx]
Checks, on the given file: no [..] input slots remain; the synthetic-data
disclaimer (VN + EN) is present on every flagged slide; no forbidden claim
phrase appears in slide text or speaker notes.
Exit 0 = ready; exit 1 = problems (listed); exit 2 = file unreadable.
"""
import sys

from pptx import Presentation
from pptx.exc import PackageNotFoundError

from deck.compliance import FORBIDDEN, SLOT_RE
from deck.content import DISCLAIMER_SLIDES
from deck.theme import DISCLAIMER_EN, DISCLAIMER_VN


def slide_texts(slide):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    texts.append(cell.text)
    return texts


def check(path):
    prs = Presentation(path)
    problems = []
    for idx, slide in enumerate(prs.slides, start=1):
        joined = "\n".join(slide_texts(slide))
        for token in SLOT_RE.findall(joined):
            problems.append(f"slide {idx}: unfilled slot {token}")
        if idx in DISCLAIMER_SLIDES:
            if DISCLAIMER_VN not in joined:
                problems.append(f"slide {idx}: VN disclaimer missing")
            if DISCLAIMER_EN not in joined:
                problems.append(f"slide {idx}: EN disclaimer missing")
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
        lowered = (joined + "\n" + notes).lower()
        for phrase in FORBIDDEN:
            if phrase.lower() in lowered:
                problems.append(f"slide {idx}: forbidden phrase '{phrase}'")
    return problems


def main():
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    path = sys.argv[1] if len(sys.argv) > 1 else "deck/output/deck.pptx"
    try:
        problems = check(path)
    except (FileNotFoundError, PackageNotFoundError):
        print(f"ERROR: cannot open '{path}' — check the file path.")
        return 2
    if not problems:
        print(f"OK: {path} passes the finalization gate.")
        return 0
    print(f"NOT FINAL: {len(problems)} problem(s) in {path}:")
    for problem in problems:
        print(f"  {problem}")
    return 1


if __name__ == "__main__":
    sys.exit(main())
