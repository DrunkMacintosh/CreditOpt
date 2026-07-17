import re

import pytest
from pptx import Presentation

from deck.content import DISCLAIMER_SLIDES
from deck.theme import DISCLAIMER_VN, DISCLAIMER_EN

FORBIDDEN = [
    "đầu tiên của SHB",
    "được SHB phê duyệt",
    "SHB đã phê duyệt",
    "SHB chứng thực",
    "production-ready",
    "sẵn sàng production",
    "đã được chứng nhận bảo mật",
]


@pytest.fixture(scope="session")
def prs():
    from deck.build import build_deck
    return Presentation(build_deck("deck/output/test_deck.pptx"))


def slide_text(slide):
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    parts.append(cell.text)
    return "\n".join(parts)


def test_disclaimer_on_required_slides(prs):
    for idx, slide in enumerate(prs.slides, start=1):
        if idx in DISCLAIMER_SLIDES:
            text = slide_text(slide)
            assert DISCLAIMER_VN in text, f"slide {idx}: VN disclaimer missing"
            assert DISCLAIMER_EN in text, f"slide {idx}: EN disclaimer missing"


def test_no_forbidden_claims(prs):
    for idx, slide in enumerate(prs.slides, start=1):
        text = slide_text(slide).lower()
        for phrase in FORBIDDEN:
            assert phrase.lower() not in text, f"slide {idx}: forbidden '{phrase}'"


def test_input_slots_still_present_until_filled(prs):
    # Guards against someone silently deleting the slots instead of filling them.
    slides = list(prs.slides)
    for idx in (13, 15, 17):
        assert re.search(r"\[[^\]]+\]", slide_text(slides[idx - 1])), \
            f"slide {idx}: expected input slots"
